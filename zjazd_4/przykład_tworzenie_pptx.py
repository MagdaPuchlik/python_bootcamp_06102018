from pptx import Presentation

prs = Presentation()

slide_layout = prs.slide_layouts[1] #wybieramy z ajkieś listy layout sladjy

slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "jakiś tekst"

tf = body_shape.text_frame
tf.text = "zawartość tekst frame"

p = tf.add_paragraph()
p.text = "Kobiety"
p.level = 1

p = tf.add_paragraph()
p.text = "Przeżyło:70"
p.level = 2

slide_layout2 = prs.slide_layouts[2] #wybieramy z ajkieś listy layout sladjy

slide = prs.slides.add_slide(slide_layout2)



prs.save("raport.pptx")
